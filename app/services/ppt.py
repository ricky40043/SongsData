from pathlib import Path
import re
import unicodedata

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from sqlalchemy import select
from sqlalchemy.orm import Session

from app.db.models import Song, SongVersion
from app.services.normalizer import clean_lines


DEFAULT_TEMPLATE_PATH = Path("app/templates/songppt-template.pptx")
MAX_DISPLAY_LINES_PER_SLIDE = 2
MAX_DISPLAY_WIDTH_PER_LINE = 16.0
MAX_DISPLAY_WIDTH_PER_SLIDE = MAX_DISPLAY_WIDTH_PER_LINE * MAX_DISPLAY_LINES_PER_SLIDE
LYRIC_FONT_SIZE_PT = 48
LYRIC_SPLIT_RE = re.compile(r"([，,；;。．.!！?？、])")
DISPLAY_PUNCTUATION_RE = re.compile(r"[，,；;。．.!！?？、]")
HARD_BREAK_PUNCTUATION = {"；", ";", "。", "．", ".", "!", "！", "?", "？"}
SECTION_LABEL_RE = re.compile(r"^(?:副歌|主歌|橋段|間奏|前奏|尾奏|結尾|簡版|chorus|verse|bridge)\s*[:：]?\s*", re.I)
METADATA_KEYWORDS = (
    "詩歌",
    "歌集",
    "專輯",
    "讚美之泉",
    "小羊",
    "新靈糧",
    "青年聖歌",
    "齊唱短歌",
    "敬拜讚美",
    "我心旋律",
)
UNSAFE_FILENAME_RE = re.compile(r'[\\/:*?"<>|\r\n\t]+')


def generate_song_pptx(
    db: Session,
    song_id: int,
    output_path: str | Path,
    template_path: str | Path = DEFAULT_TEMPLATE_PATH,
) -> Path:
    song = db.get(Song, song_id)
    if not song:
        raise ValueError(f"Song not found: {song_id}")

    version = db.scalar(select(SongVersion).where(SongVersion.song_id == song_id))
    if not version or not version.raw_lyrics.strip():
        raise ValueError(f"No lyrics found for song: {song_id}")

    lyric_slides = build_songppt_slide_payloads(version.raw_lyrics)
    if not lyric_slides:
        raise ValueError(f"No lyric slides found for song: {song_id}")

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)

    template = Path(template_path)
    prs = Presentation(template) if template.exists() else Presentation()
    _remove_all_slides(prs)
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    for slide_model in lyric_slides:
        slide = prs.slides.add_slide(blank_layout)
        _clear_shapes(slide)
        _set_black_background(slide)
        _add_songppt_text(slide, slide_model["lines"], slide_model.get("repeat", 1))

    prs.save(output)
    return output


def song_pptx_filename(title: str, song_id: int | None = None) -> str:
    cleaned = UNSAFE_FILENAME_RE.sub("_", title).strip(" ._")
    cleaned = re.sub(r"\s+", " ", cleaned)
    if not cleaned:
        cleaned = f"song-{song_id}" if song_id is not None else "song"
    return f"{cleaned[:80]}.pptx"


def build_songppt_slide_payloads(raw_lyrics: str) -> list[dict]:
    unmerged = []
    for lines in _prepare_display_slides(raw_lyrics):
        if lines:
            unmerged.append({"lines": lines})

    merged: list[dict] = []
    for slide in unmerged:
        if merged and merged[-1]["lines"] == slide["lines"]:
            merged[-1]["repeat"] += 1
        else:
            merged.append({"lines": slide["lines"], "repeat": 1})
    return merged


def _prepare_display_slides(raw_lyrics: str) -> list[list[str]]:
    slides: list[list[str]] = []
    pending_single_lines: list[str] = []
    source_lines = clean_lines(raw_lyrics)
    for idx, raw_line in enumerate(source_lines):
        if _is_metadata_line(raw_line, idx):
            continue
        for logical_part in _split_logical_parts(raw_line):
            wrapped = _wrap_display_line(logical_part, MAX_DISPLAY_WIDTH_PER_LINE)
            if len(wrapped) == 1:
                pending_single_lines.append(wrapped[0])
                if len(pending_single_lines) == MAX_DISPLAY_LINES_PER_SLIDE:
                    slides.append(pending_single_lines)
                    pending_single_lines = []
                continue

            if pending_single_lines:
                slides.append(pending_single_lines)
                pending_single_lines = []
            slides.extend(_wrap_chunks_to_slides(wrapped))
    if pending_single_lines:
        slides.append(pending_single_lines)
    return slides


def _is_metadata_line(line: str, idx: int) -> bool:
    text = line.strip()
    if not text:
        return True
    if idx <= 2 and any(keyword in text for keyword in METADATA_KEYWORDS):
        return True
    if idx <= 2 and text.startswith(("(", "（")) and text.endswith((")", "）")):
        return True
    return False


def _split_logical_parts(line: str) -> list[str]:
    normalized = line.replace("｜", "|").replace("\u000b", "|")
    parts: list[str] = []
    for pipe_part in normalized.split("|"):
        cleaned = _clean_display_text(pipe_part)
        if not cleaned:
            continue
        parts.append(cleaned)
    return parts


def _clean_display_text(text: str) -> str:
    cleaned = text.strip()
    cleaned = re.sub(r"([哦喔噢])\s*[!！]\s*(主聖靈|聖靈)", r"\1\2", cleaned)
    cleaned = SECTION_LABEL_RE.sub("", cleaned).strip()
    if cleaned.startswith(("(", "（")) and cleaned.endswith((")", "）")):
        return ""
    return cleaned


def _split_by_punctuation(text: str) -> list[str]:
    tokens = LYRIC_SPLIT_RE.split(text)
    groups: list[list[str]] = []
    clauses: list[str] = []
    current = ""
    for token in tokens:
        if not token:
            continue
        if LYRIC_SPLIT_RE.fullmatch(token):
            if current.strip():
                clauses.append(f"{current.strip()}{token}")
                current = ""
            if token in HARD_BREAK_PUNCTUATION and clauses:
                groups.append(clauses)
                clauses = []
            continue
        current += token
    if current.strip():
        clauses.append(current.strip())

    if clauses:
        groups.append(clauses)

    if not groups:
        groups = [[text.strip()]]

    packed: list[str] = []
    for group in groups:
        current = ""
        for clause in group:
            if not clause:
                continue
            candidate = (
                f"{current} {clause}".strip()
                if _contains_ascii_word(current + clause)
                else f"{current}{clause}"
            )
            if current and _display_width(candidate) > MAX_DISPLAY_WIDTH_PER_LINE:
                packed.append(current)
                current = clause
            else:
                current = candidate
        if current:
            packed.append(current)
    return packed


def _strip_display_punctuation(text: str) -> str:
    return text.strip().strip("，,；;。．.!！?？、 ")


def _wrap_display_line(text: str, max_width: float) -> list[str]:
    if _display_width(text) <= max_width:
        cleaned = _remove_display_punctuation(text)
        return [cleaned] if cleaned else []

    chunks: list[str] = []
    remaining = text.strip()
    while remaining:
        if _display_width(remaining) <= max_width:
            chunks.append(_remove_display_punctuation(remaining))
            break

        break_idx = _best_wrap_index(remaining, max_width)
        if break_idx <= 0:
            break_idx = _hard_wrap_index(remaining, max_width)
        chunks.append(_remove_display_punctuation(remaining[:break_idx]))
        remaining = remaining[break_idx:].strip()
    return [chunk for chunk in chunks if chunk]


def _wrap_chunks_to_slides(chunks: list[str]) -> list[list[str]]:
    if len(chunks) <= MAX_DISPLAY_LINES_PER_SLIDE:
        return [chunks]

    slides: list[list[str]] = []
    for idx in range(0, len(chunks), MAX_DISPLAY_LINES_PER_SLIDE):
        slides.append(chunks[idx : idx + MAX_DISPLAY_LINES_PER_SLIDE])
    return slides


def _best_wrap_index(text: str, max_width: float) -> int:
    best = -1
    for idx, char in enumerate(text):
        if _display_width(text[: idx + 1]) > max_width:
            break
        if char.isspace() or char in "，,；;。．.!！?？、":
            best = idx + 1
    return best


def _hard_wrap_index(text: str, max_width: float) -> int:
    current = ""
    for idx, char in enumerate(text):
        candidate = current + char
        if current and _display_width(candidate) > max_width:
            return idx
        current = candidate
    return len(text)


def _remove_display_punctuation(text: str) -> str:
    return DISPLAY_PUNCTUATION_RE.sub("", text).strip()


def _display_width(text: str) -> float:
    width = 0.0
    for char in text:
        if char.isspace():
            width += 0.35
        elif unicodedata.east_asian_width(char) in {"F", "W", "A"}:
            width += 1.0
        else:
            width += 0.55
    return width


def _contains_ascii_word(text: str) -> bool:
    return bool(re.search(r"[A-Za-z]", text))


def _remove_all_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def _clear_shapes(slide) -> None:
    shape_tree = slide.shapes._spTree  # noqa: SLF001 - python-pptx has no public clear API.
    for shape in list(slide.shapes):
        shape_tree.remove(shape._element)  # noqa: SLF001


def _set_black_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)


def _add_songppt_text(slide, lines: list[str], repeat: int) -> None:
    textbox = slide.shapes.add_textbox(Inches(0.45), Inches(0.05), Inches(12.45), Inches(2.55))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    display_lines = lines[:MAX_DISPLAY_LINES_PER_SLIDE]
    if repeat > 1:
        display_lines.append(f"(x{repeat})")

    for idx, text in enumerate(display_lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.space_after = Pt(0)
        paragraph.space_before = Pt(0)
        paragraph.line_spacing = 0.9
        run = paragraph.add_run()
        run.text = text
        run.font.name = "源泉圓體 B"
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 0)
        run.font.size = _font_size_for(text, is_repeat=(repeat > 1 and idx == 2))


def _font_size_for(text: str, is_repeat: bool = False):
    if is_repeat:
        return Pt(38)
    return Pt(LYRIC_FONT_SIZE_PT)
